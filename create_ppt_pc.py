from pptx import Presentation
from pptx.util import Inches, Cm, Pt
import os
from PIL import Image
from tqdm import tqdm
import natsort
from pptx.dml.color import RGBColor
import math
from pptx.enum.text import PP_ALIGN

root_path = 'target_pc'
folder_list = os.listdir(root_path)

for folder in folder_list:
    print("대상 폴더명 : {folder}".format(folder=folder))
    file_list = os.listdir(root_path + '/' + folder)
    file_list = natsort.natsorted(file_list)
    img_num = 1
    top = Cm(0)
    prs = Presentation()  # 파워포인트 객체 선언
    prs.slide_height = Inches(9)
    prs.slide_width = Inches(16)

    for file in tqdm(file_list, desc="Create PPTX"):
        img_path = root_path + '/' + folder + '/' + file
        image = Image.open(img_path)
        img_width, img_height = image.size

        left = Cm(0)
        top = Cm(0)
        # text_top = Cm(20.19)
        text_top = Cm(17.12)
        # next_left = Cm(26.01)
        next_left = Cm(26.24)

# 기존: 1496
        if img_height > 1250:
            # start_x, start_y, start_x + width, start_y + height
            crop_img = []
            crop_cnt = 1
            crop_max = math.ceil(img_height/1250.0)
            for i in range(0, img_height, 1250):
                box = (0, i, img_width, i + 1250)
                crop_image = image.crop(box)
                crop_image.save('tmp_pc_img.png')
                blank_slide_layout = prs.slide_layouts[6]  # 6 : 제목/내용이 없는 '빈' 슬라이드
                slide = prs.slides.add_slide(blank_slide_layout)

                #pic = slide.shapes.add_picture('tmp_pc_img.png', left, top, width=Cm(26.01))
                pic = slide.shapes.add_picture('tmp_pc_img.png', left, top, width=Cm(26.24))

#                txbox = slide.shapes.add_textbox(left, text_top, width=Cm(26.01), height=Cm(2.68))
                txbox = slide.shapes.add_textbox(left, text_top, width=Cm(26.24), height=Cm(2.68))
                tf = txbox.text_frame
                tf.word_wrap = True

                fill = txbox.fill
                fill.solid()

                fill.fore_color.rgb = RGBColor(38, 50, 70)

                p = tf.add_paragraph()
                p.text = file
                p.font.size = Pt(12)
                p.font.name = '맑은 고딕'
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)

                if crop_cnt < crop_max:
                    txbox = slide.shapes.add_textbox(next_left, text_top, width=Cm(14.63), height=Cm(2.68))
                    tf = txbox.text_frame
                    tf.word_wrap = True

                    fill = txbox.fill
                    fill.solid()

                    fill.fore_color.rgb = RGBColor(255, 0, 0)

                    p = tf.add_paragraph()
                    p.text = "NEXT"
                    p.font.size = Pt(18)
                    p.font.name = '맑은 고딕'
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(255, 255, 255)
                    p.alignment = PP_ALIGN.CENTER

                img_num += 1
                crop_cnt += 1

        else:
            blank_slide_layout = prs.slide_layouts[6]  # 6 : 제목/내용이 없는 '빈' 슬라이드
            slide = prs.slides.add_slide(blank_slide_layout)

            image.save('tmp_img.png')
            pic = slide.shapes.add_picture('tmp_img.png', left, top, width=Cm(26.01))
            txbox = slide.shapes.add_textbox(left, text_top, width=Cm(26.01), height=Cm(2.68))
            tf = txbox.text_frame
            tf.word_wrap = True

            fill = txbox.fill
            fill.solid()

            fill.fore_color.rgb = RGBColor(38, 50, 70)

            p = tf.add_paragraph()
            p.text = file
            p.font.size = Pt(12)
            p.font.name = '맑은 고딕'
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)

            img_num += 1

    prs.save('{file_name}.pptx'.format(file_name=folder))
