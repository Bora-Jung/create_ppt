"""
나는 나는 저팔계

"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from pptx.dml.color import RGBColor


# set variables(초기값)
PPT_IMG_WIDTH = Cm(7.1)
PPT_IMG_HEIGHT = Cm(20)

SLIDE_HEIGHT = 9
SLIDE_WIDTH = 16

CM_TO_PIXEL = 96
CUT_HEIGHT_CM = 17.3

IMG_Y = 0

FOLDER_PATH = 'target/@UHDC_Part5_PC_MW_상담신청_0831'


class BoraPptCreator:
    """
    class로 만들어 보기 (나는 class의 document다!)
    """
    
    def __init__(self,
                 folder_path: str = FOLDER_PATH,
                 slide_height: str = SLIDE_HEIGHT,
                 slide_width: str = SLIDE_WIDTH):
        """
        init parameter 설명!

        :param folder_path:
        :param slide_height:
        :param slide_width:
        """

        # start ppt
        self.prs = Presentation()

        # set property on ppt
        self.prs.slide_height = Inches(slide_height)
        self.prs.slide_width = Inches(slide_width)

        # initialization
        self.folder_path = folder_path
        self.slide = None

    # x 좌표
    @staticmethod
    def find_x_point(cnt: int) -> Cm:
        if cnt % 4 == 1:
            x = 0
        elif cnt % 4 == 2:
            x = Cm(11.18)
        elif cnt % 4 == 3:
            x = Cm(22.36)
        elif cnt % 4 == 0:
            x = Cm(33.54)
        else:
            raise AssertionError("말도안댐")

        return x

    # 슬라이드에 넣기
    def make_slide(self, x_location: int, file_name: str):

        # slide 생성
        self.slide.shapes.add_picture('tmp.png', x_location, IMG_Y, PPT_IMG_WIDTH, PPT_IMG_HEIGHT)

        # textbox 생성
        tb = self.slide.shapes.add_textbox(x_location, PPT_IMG_HEIGHT, PPT_IMG_WIDTH, Inches(1))
        tf = tb.text_frame
        p = tf.paragraphs[0]

        # ppt
        run = p.add_run()
        run.text = file_name

        # font 만들기 & 스타일 적용
        font = run.font
        font.name = '맑은 고딕'
        font.size = Pt(11)

        # textbox 색칠하기
        fill = tb.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(217, 217, 217)

    def create(self):

        sum_cut = 0

        for i, file in enumerate(os.listdir(self.folder_path)):

            # open an image
            img = Image.open(os.path.join(self.folder_path, file))

            # get size
            width_px, height_px = img.size

            # cm -> px
            width_new_px = 7 * CM_TO_PIXEL
            height_new_px = height_px * width_new_px / width_px

            # image resizing
            img_new = img.resize((int(width_new_px), int(height_new_px)))

            # px -> cm
            cut_height_px = CUT_HEIGHT_CM * CM_TO_PIXEL

            # cutting size
            cut_cnt = int((height_new_px / cut_height_px) + 1)

            for cut in range(cut_cnt):
                top = cut * cut_height_px
                bottom = cut_height_px * (cut + 1)
                cut_size = (0, top, width_new_px, bottom)
                img_crop = img_new.crop(cut_size)
                img_crop.save('tmp.png')

                # file_name
                file_name = file.split('.png')[0] if cut == 0 else ''

                sum_cut += 1
                img_x = self.find_x_point(sum_cut)

                if sum_cut % 4 == 1:
                    slide_layout = self.prs.slide_layouts[6]
                    self.slide = self.prs.slides.add_slide(slide_layout)

                # slide 만들깅
                self.make_slide(x_location=img_x, file_name=file_name)

            # find path
            crop_folder = self.folder_path.split('/')[1]

            # save
            self.prs.save(f"{crop_folder}.pptx")


if __name__ == '__main__':
    import argparse

    parser = argparse.ArgumentParser()
    parser.add_argument('--fp',
                        type=str)
    parser.add_argument('--sh',
                        type=int)
    parser.add_argument('--sw',
                        type=int)

    parsed_args = parser.parse_args()

    folder_path = parsed_args.fp
    slide_height = parsed_args.sh
    slide_width = parsed_args.sw

    print(folder_path, slide_height, slide_width)

    bora_creator = BoraPptCreator(folder_path=folder_path,
                                  slide_height=slide_height,
                                  slide_width=slide_width,)
    bora_creator.create()
