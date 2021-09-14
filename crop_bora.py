import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from pptx.dml.color import RGBColor


# set variables(초기값)
PPT_IMG_WIDTH = Cm(7.1)
PPT_IMG_HEIGHT = Cm(20)

SLIDE_HEIGHT = Inches(9)
SLIDE_WIDTH = Inches(16)

CM_TO_PIXEL = 96
CUT_HEIGHT_CM = 17.3

IMG_Y = 0

FOLDER_PATH = 'target/@UHDC_Part5_PC_MW_상담신청_0831'


# x 좌표
def x_point(cnt: int) -> Cm:
    if cnt % 4 == 1:
        x = 0
    elif cnt % 4 == 2:
        x = Cm(11.18)
    elif cnt % 4 == 3:
        x = Cm(22.36)
    elif cnt % 4 == 0:
        x = Cm(33.54)
    else:
        raise AssertionError("말도안댐")

    return x


# 슬라이드에 넣기
def make_slide(x_location: int):

    # slide 생성
    slide.shapes.add_picture('tmp.png', x_location, IMG_Y, PPT_IMG_WIDTH, PPT_IMG_HEIGHT)

    # textbox 생성
    tb = slide.shapes.add_textbox(x_location, PPT_IMG_HEIGHT, PPT_IMG_WIDTH, Inches(1))
    tf = tb.text_frame
    p = tf.paragraphs[0]

    # ppt
    run = p.add_run()
    run.text = file_name

    # font 만들기 & 스타일 적용
    font = run.font
    font.name = '맑은 고딕'
    font.size = Pt(11)

    # textbox 색칠하기
    fill = tb.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(217, 217, 217)


if __name__ == '__main__':
    # start ppt
    prs = Presentation()

    # set property on ppt
    prs.slide_height = SLIDE_HEIGHT
    prs.slide_width = SLIDE_WIDTH

    # initialize loop
    sum_cut = 0

    for i, file in enumerate(os.listdir(FOLDER_PATH)):

        # open an image
        img = Image.open(os.path.join(FOLDER_PATH, file))

        # get size
        width_px, height_px = img.size

        # cm -> px
        width_new_px = 7 * CM_TO_PIXEL
        height_new_px = height_px * width_new_px / width_px

        # image resizing
        img_new = img.resize((int(width_new_px), int(height_new_px)))

        # px -> cm
        cut_height_px = CUT_HEIGHT_CM * CM_TO_PIXEL

        # cutting size
        cut_cnt = int((height_new_px / cut_height_px) + 1)

        for cut in range(cut_cnt):
            top = cut * cut_height_px
            bottom = cut_height_px * (cut + 1)
            cut_size = (0, top, width_new_px, bottom)
            img_crop = img_new.crop(cut_size)
            img_crop.save('tmp.png')

            # file_name
            file_name = file.split('.png')[0] if cut == 0 else ''

            sum_cut += 1
            img_x = x_point(sum_cut)

            if sum_cut % 4 == 1:
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)

            # slide 만들깅
            make_slide(img_x)

        # find path
        crop_folder = FOLDER_PATH.split('/')[1]

        # save
        prs.save(f"{crop_folder}.pptx")

