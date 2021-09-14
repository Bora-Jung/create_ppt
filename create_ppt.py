from pptx import Presentation   # 라이브러리
from pptx.util import Inches, Cm, Pt    # 사진, 표등을 그리기 위해
import os
from PIL import Image
from tqdm import tqdm
import natsort
from pptx.dml.color import RGBColor

# pip install python-pptx
# pip install pillow
# pip install tqdm
# pip install natsort

def location_left(img_num):
    if img_num % 4 == 1:
        left = Cm(0)
    elif img_num % 4 == 2:
        left = Cm(10.57)
    elif img_num % 4 == 3:
        left = Cm(21.13)
    elif img_num % 4 == 0:
        left = Cm(31.7)

    return left


root_path = 'target'
folder_list = os.listdir(root_path)

for folder in folder_list:
    print("대상 폴더명 : {folder}".format(folder=folder))
    file_list = os.listdir(root_path + '/' + folder)
    file_list = natsort.natsorted(file_list)
    img_num = 1
    top = Cm(0)
    prs = Presentation()  # 파워포인트 객체 선언
    prs.slide_height = Inches(9)
    prs.slide_width = Inches(16)
    for file in tqdm(file_list, desc="Create PPTX"):
        img_path = root_path + '/' + folder + '/' + file

        image = Image.open(img_path)
        img_width, img_height = image.size

        if img_width > img_height:
            image = image.transpose(Image.ROTATE_90)
            img_width, img_height = image.size

        new_width = int(337)
        new_height = int(new_width * img_height / img_width)
        img = image.resize((new_width, new_height), Image.ANTIALIAS)

        # 이미지 분할
        if new_height > 832:
            # start_x, start_y, start_x + width, start_y + height
            crop_img = []
            crop_cnt = 1
            for i in range(0, new_height, 763):
                box = (0, i, new_width, i + 763)
                crop_image = img.crop(box)
                crop_image.save('tmp_img.png')
                if img_num % 4 == 1:
                    blank_slide_layout = prs.slide_layouts[6]  # 6 : 제목/내용이 없는 '빈' 슬라이드
                    slide = prs.slides.add_slide(blank_slide_layout)

                left = location_left(img_num)
                # with io.BytesIO() as output:
                #     crop_image.save(output, format="PNG")
                    # pic = slides.add_slide(output, left, top)
                pic = slide.shapes.add_picture('tmp_img.png', left, top, width=Cm(8.94))

                text_top = Cm(20.19)
                txbox = slide.shapes.add_textbox(left, text_top, width=Cm(8.94), height=Cm(2.68))
                tf = txbox.text_frame
                tf.word_wrap = True

                fill = txbox.fill
                fill.solid()
                # fill.patterned()
                fill.fore_color.rgb = RGBColor(38, 50, 70)


                p = tf.add_paragraph()
                p.text = file + ' ({cnt})'.format(cnt=crop_cnt)
                p.font.size = Pt(12)
                p.font.name = '맑은 고딕'
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)

                img_num += 1
                crop_cnt += 1

        else:
            if img_num % 4 == 1:
                blank_slide_layout = prs.slide_layouts[6]  # 6 : 제목/내용이 없는 '빈' 슬라이드
                slide = prs.slides.add_slide(blank_slide_layout)

            left = location_left(img_num)
            # with io.BytesIO() as output:
            #     img.save(output, format="PNG")
                # pic = slides.add_slide(output, left, top)
            img.save('tmp_img.png')
            pic = slide.shapes.add_picture('tmp_img.png', left, top, width=Cm(8.94))
            # pic = slide.shapes.add_picture(img, left, 0, width=new_width, height=new_height)

            text_top = Cm(20.19)
            txbox = slide.shapes.add_textbox(left, text_top, width=Cm(8.94), height=Cm(2.68))
            tf = txbox.text_frame
            tf.word_wrap = True

            fill = txbox.fill
            fill.solid()
            # fill.patterned()
            fill.fore_color.rgb = RGBColor(38, 50, 70)

            p = tf.add_paragraph()
            p.text = file
            p.font.size = Pt(12)
            p.font.name = '맑은 고딕'
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)

            img_num += 1

    prs.save('{file_name}.pptx'.format(file_name=folder))
